#!/usr/bin/env python3
"""
Education Center Management System - Professional Presentation Generator
Generates a comprehensive PowerPoint presentation analyzing the project.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Color scheme
PRIMARY_COLOR = RGBColor(0x1e, 0x40, 0xaf)  # Blue
SECONDARY_COLOR = RGBColor(0x06, 0x3b, 0x5e)  # Dark blue
ACCENT_COLOR = RGBColor(0x10, 0xb9, 0x81)  # Green
TEXT_COLOR = RGBColor(0x37, 0x47, 0x4f)  # Dark gray
LIGHT_BG = RGBColor(0xf8, 0xfa, 0xfc)  # Light gray background

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

    return slide

def create_content_slide(prs, title, content_points, layout=1):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True

    if layout == 1:  # Title and Content layout
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()

        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0

    return slide

def create_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Create a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Left column
    left_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4.5), Inches(5))
    left_tf = left_box.text_frame
    left_tf.clear()

    # Left title
    p = left_tf.add_paragraph()
    p.text = left_title
    p.font.size = Pt(24)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True

    for point in left_points:
        p = left_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6), Inches(1.5), Inches(4.5), Inches(5))
    right_tf = right_box.text_frame
    right_tf.clear()

    # Right title
    p = right_tf.add_paragraph()
    p.text = right_title
    p.font.size = Pt(24)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True

    for point in right_points:
        p = right_tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR

    return slide

def create_architecture_slide(prs, title, components):
    """Create an architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Create architecture boxes
    y_pos = Inches(1.5)
    for i, component in enumerate(components):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(2 + (i % 2) * 4),
            y_pos,
            Inches(3),
            Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_COLOR if i < 2 else PRIMARY_COLOR
        box.line.color.rgb = SECONDARY_COLOR

        tf = box.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = component
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if i % 2 == 1:
            y_pos += Inches(1.5)

    return slide

def create_features_grid_slide(prs, title, features):
    """Create a features grid slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Create 2x3 grid of features
    positions = [
        (Inches(1), Inches(1.5), Inches(3.5), Inches(1.5)),
        (Inches(5.5), Inches(1.5), Inches(3.5), Inches(1.5)),
        (Inches(1), Inches(3.5), Inches(3.5), Inches(1.5)),
        (Inches(5.5), Inches(3.5), Inches(3.5), Inches(1.5)),
        (Inches(1), Inches(5.5), Inches(3.5), Inches(1.5)),
        (Inches(5.5), Inches(5.5), Inches(3.5), Inches(1.5)),
    ]

    for i, feature in enumerate(features[:6]):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            positions[i][0], positions[i][1], positions[i][2], positions[i][3]
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = PRIMARY_COLOR

        tf = box.text_frame
        tf.clear()
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_tech_stack_slide(prs, title, backend_tech, frontend_tech, database_tech, deployment_tech):
    """Create technology stack slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Backend column
    backend_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3), Inches(5))
    backend_tf = backend_box.text_frame
    backend_tf.clear()

    p = backend_tf.add_paragraph()
    p.text = "Backend"
    p.font.size = Pt(24)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True

    for tech in backend_tech:
        p = backend_tf.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR

    # Frontend column
    frontend_box = slide.shapes.add_textbox(Inches(4), Inches(1.5), Inches(3), Inches(5))
    frontend_tf = frontend_box.text_frame
    frontend_tf.clear()

    p = frontend_tf.add_paragraph()
    p.text = "Frontend"
    p.font.size = Pt(24)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True

    for tech in frontend_tech:
        p = frontend_tf.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR

    # Database column
    database_box = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(2.5), Inches(2.5))
    database_tf = database_box.text_frame
    database_tf.clear()

    p = database_tf.add_paragraph()
    p.text = "Database"
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True

    for tech in database_tech:
        p = database_tf.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR

    # Deployment column
    deployment_box = slide.shapes.add_textbox(Inches(7.5), Inches(4.5), Inches(2.5), Inches(2.5))
    deployment_tf = deployment_box.text_frame
    deployment_tf.clear()

    p = deployment_tf.add_paragraph()
    p.text = "Deployment"
    p.font.size = Pt(20)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True

    for tech in deployment_tech:
        p = deployment_tf.add_paragraph()
        p.text = f"• {tech}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR

    return slide

def create_conclusion_slide(prs, title, points):
    """Create conclusion slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_shape.text_frame.paragraphs[0].font.bold = True

    content_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
    tf = content_box.text_frame
    tf.clear()

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    """Generate the complete presentation"""
    prs = Presentation()

    # Slide 1: Title
    create_title_slide(
        prs,
        "نظام إدارة مركز تعليمي",
        "Education Center Management System\nComprehensive Learning Platform"
    )

    # Slide 2: Project Overview
    create_content_slide(
        prs,
        "نظرة عامة على المشروع",
        [
            "منصة تعليمية متكاملة تربط بين الطلاب والمعلمين",
            "إدارة شاملة للعملية التعليمية من التسجيل إلى التخرج",
            "دعم كامل للتعليم التقليدي والإلكتروني",
            "حلول تقنية متقدمة لإدارة المحتوى والتقييم"
        ]
    )

    # Slide 3: Problem Statement
    create_content_slide(
        prs,
        "المشكلة التي يحلها النظام",
        [
            "صعوبة إدارة العمليات التعليمية يدوياً",
            "عدم وجود منصة موحدة للتواصل والمتابعة",
            "تحديات في تتبع التقدم الأكاديمي للطلاب",
            "مشاكل في إدارة المدفوعات والتسجيل",
            "صعوبة في تنظيم الجدول الدراسي والحضور"
        ]
    )

    # Slide 4: Core Features
    create_features_grid_slide(
        prs,
        "المميزات الأساسية",
        [
            "🔐 المصادقة والصلاحيات",
            "📚 إدارة الكورسات",
            "📝 نظام الالتحاق والدفع",
            "📺 المحتوى التعليمي",
            "📊 تتبع التقدم",
            "📅 إدارة الجدول الدراسي",
            "🏛️ إدارة القاعات",
            "📝 نظام الحضور",
            "🧠 الاختبارات المتقدمة",
            "📋 نظام الواجبات",
            "💬 نظام المراسلة",
            "👥 إدارة المستخدمين"
        ]
    )

    # Slide 5: User Roles
    create_two_column_slide(
        prs,
        "أدوار المستخدمين",
        "الطالب (Student)",
        [
            "تصفح والبحث عن الكورسات",
            "طلب الالتحاق في الكورسات",
            "متابعة الدروس والمحتوى",
            "إجراء الاختبارات والواجبات",
            "تتبع التقدم الشخصي",
            "عرض الجدول الدراسي"
        ],
        "المعلم (Teacher)",
        [
            "إنشاء وإدارة الكورسات",
            "إضافة المحتوى التعليمي",
            "قبول/رفض طلبات الطلاب",
            "إنشاء الاختبارات والواجبات",
            "تصحيح وتقييم الأعمال",
            "متابعة أداء الطلاب"
        ]
    )

    # Slide 6: Payment System
    create_content_slide(
        prs,
        "نظام الدفع الإلكتروني",
        [
            "🔄 دورة حياة الدفع:",
            "  1. طلب الالتحاق → 2. المراجعة → 3. الدفع → 4. التفعيل",
            "💳 تكامل مع بوابة Paymob",
            "💳 دعم البطاقات البنكية (Visa/Mastercard)",
            "📱 دعم المحافظ الإلكترونية",
            "⚡ معالجة فورية وتفعيل تلقائي",
            "📈 دعم الدفع الشهري والمرة واحدة"
        ]
    )

    # Slide 7: Content Management
    create_two_column_slide(
        prs,
        "إدارة المحتوى التعليمي",
        "أنواع المحتوى",
        [
            "🎬 فيديوهات تعليمية",
            "📄 ملفات PDF ومذكرات",
            "🖼️ صور توضيحية",
            "📝 نصوص وشروحات مكتوبة"
        ],
        "المميزات التقنية",
        [
            "⏸️ مشغل فيديو مخصص",
            "💾 حفظ تلقائي للتقدم",
            "📊 تتبع المشاهدة الفعلية",
            "🔒 منع التخطي (اختياري)",
            "📱 دعم التصميم المتجاوب"
        ]
    )

    # Slide 8: Assessment System
    create_two_column_slide(
        prs,
        "نظام التقييم والاختبارات",
        "الاختبارات (Quizzes)",
        [
            "📝 اختبارات قصيرة ونهائية",
            "⏱️ مؤقت تنازلي",
            "🔒 منع الغش والتخطي",
            "📊 تصحيح آلي فوري",
            "👁️ مراقبة حية للمعلم"
        ],
        "الواجبات (Assignments)",
        [
            "📤 رفع الملفات (PDF)",
            "📅 مواعيد التسليم",
            "👨‍🏫 واجهة التصحيح للمعلم",
            "💬 نظام الملاحظات والتقييم",
            "📈 تتبع حالة التسليم"
        ]
    )

    # Slide 9: Architecture
    create_architecture_slide(
        prs,
        "المعمارية التقنية",
        [
            "React Frontend",
            "Vite Build Tool",
            "NestJS Backend",
            "MySQL Database",
            "Prisma ORM",
            "Docker Containers",
            "Nginx Web Server",
            "Paymob Payment Gateway"
        ]
    )

    # Slide 10: Technology Stack
    create_tech_stack_slide(
        prs,
        "التقنيات المستخدمة",
        [
            "NestJS Framework",
            "Node.js Runtime",
            "TypeScript",
            "Prisma ORM",
            "JWT Authentication",
            "Passport.js",
            "Swagger API Docs",
            "Socket.io (Real-time)"
        ],
        [
            "React 18",
            "Vite Build Tool",
            "React Router",
            "Axios HTTP Client",
            "Framer Motion",
            "Recharts (Charts)",
            "Lucide Icons",
            "SweetAlert2"
        ],
        [
            "MySQL 8.0",
            "phpMyAdmin",
            "Docker Volumes",
            "Database Migrations"
        ],
        [
            "Docker & Compose",
            "Nginx Reverse Proxy",
            "Multi-stage Builds",
            "Environment Config"
        ]
    )

    # Slide 11: Database Schema
    create_content_slide(
        prs,
        "قاعدة البيانات",
        [
            "👤 User: المستخدمين وأدوارهم",
            "📚 Course: الكورسات والمحتوى",
            "📖 Chapter: الفصول والدروس",
            "📝 Enrollment: التسجيلات والاشتراكات",
            "💰 Payment: المدفوعات والمعاملات",
            "📊 ChapterProgress: تتبع التقدم",
            "📋 Quiz/Question/Option: نظام الاختبارات",
            "📄 Assignment: الواجبات والتسليمات",
            "🏫 Room: إدارة القاعات",
            "📅 CourseSchedule: الجدول الدراسي",
            "👥 Attendance: نظام الحضور"
        ]
    )

    # Slide 12: Key Achievements
    create_content_slide(
        prs,
        "الإنجازات الرئيسية",
        [
            "✅ تكامل كامل مع بوابة الدفع Paymob",
            "✅ نظام مصادقة آمن باستخدام JWT",
            "✅ مشغل فيديو مخصص مع تتبع التقدم",
            "✅ نظام اختبارات متقدم مع مراقبة حية",
            "✅ إدارة شاملة للواجبات والتقييم",
            "✅ جدولة مرنة مع إدارة القاعات",
            "✅ نظام إشعارات ومراسلة فورية",
            "✅ تصميم متجاوب يدعم جميع الأجهزة"
        ]
    )

    # Slide 13: Deployment & DevOps
    create_content_slide(
        prs,
        "النشر والتطوير",
        [
            "🐳 Docker Containerization",
            "📦 Multi-stage Docker Builds",
            "🔄 Hot Reload Development",
            "🌐 Nginx Reverse Proxy",
            "💾 Persistent Data Volumes",
            "🔧 Environment Configuration",
            "📊 phpMyAdmin Database Admin",
            "🚀 One-command Deployment"
        ]
    )

    # Slide 14: Conclusion
    create_conclusion_slide(
        prs,
        "الخلاصة",
        [
            "منصة تعليمية متكاملة ومتقدمة",
            "حل شامل لإدارة المراكز التعليمية",
            "تقنيات حديثة ومعايير أمان عالية",
            "تجربة مستخدم ممتازة وسهلة الاستخدام",
            "جاهز للنشر والاستخدام التجاري"
        ]
    )

    # Save the presentation
    output_path = os.path.join(os.path.dirname(__file__), "Education_Center_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()